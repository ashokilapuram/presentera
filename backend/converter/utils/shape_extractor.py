"""
Phase 2: Shape Extraction from PPTX
Extracts shapes (circles, rectangles, squares, rounded rectangles, lines) with colors and properties.
"""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_FILL_TYPE, MSO_LINE_DASH_STYLE
import uuid
from typing import List, Dict, Any, Optional
from converter.schemas.slide_schema import ShapeElement


def emu_to_points(emu: int) -> float:
    """
    Convert EMU (English Metric Units) to points.
    1 inch = 914400 EMU = 72 points
    """
    return (emu / 914400) * 72


def rgb_to_hex(rgb) -> Optional[str]:
    """Convert RGB color to hex string"""
    if rgb is None:
        return None
    
    # Handle different color formats
    try:
        if hasattr(rgb, 'rgb'):
            # Check if rgb property exists and is not None
            if rgb.rgb is None:
                return None
            rgb = rgb.rgb
    except (AttributeError, TypeError):
        pass
    
    if rgb is None:
        return None
    
    # Extract RGB values
    if isinstance(rgb, int):
        # Assume it's a 32-bit integer: 0x00RRGGBB
        r = (rgb >> 16) & 0xFF
        g = (rgb >> 8) & 0xFF
        b = rgb & 0xFF
    elif hasattr(rgb, '__iter__') and len(rgb) >= 3:
        r, g, b = rgb[0], rgb[1], rgb[2]
    else:
        return None
    
    return f"#{r:02x}{g:02x}{b:02x}"


def get_shape_type(shape) -> Optional[str]:
    """
    Determine the shape type from PowerPoint shape.
    Returns: "circle", "square", "rectangle", "roundedRectangle", "line", 
             "triangle", "star", "pentagon", "hexagon", or None.
    
    Supports:
    - AUTO_SHAPE types: rectangle, circle, square, roundedRectangle, line,
                        triangle, star, pentagon, hexagon
    - FREEFORM shapes: Detects polygons by point count (3=triangle, 5=pentagon,
                       6=hexagon, 10=star)
    - Fallback: Returns "rectangle" for unrecognized shapes
    """
    try:
        # Get dimensions for comparison
        width = shape.width if hasattr(shape, 'width') else 0
        height = shape.height if hasattr(shape, 'height') else 0
        is_square = abs(width - height) < 1000  # Allow small difference (tolerance)
        
        # Get the shape's type
        if hasattr(shape, 'shape_type'):
            shape_type = shape.shape_type
            
            # Check if it's an auto shape
            if shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                if hasattr(shape, 'auto_shape_type'):
                    auto_type = shape.auto_shape_type
                    
                    # Try to import MSO_SHAPE constants
                    try:
                        from pptx.enum.shapes import MSO_SHAPE
                        
                        # Shape mapping table for new shapes
                        KNOWN_SHAPES = {
                            MSO_SHAPE.ISOSCELES_TRIANGLE: "triangle",
                            MSO_SHAPE.RIGHT_TRIANGLE: "triangle",
                            MSO_SHAPE.STAR_5_POINT: "star",
                            MSO_SHAPE.PENTAGON: "pentagon",
                            MSO_SHAPE.HEXAGON: "hexagon"
                        }
                        
                        # Check for new shapes first
                        if auto_type in KNOWN_SHAPES:
                            return KNOWN_SHAPES[auto_type]
                        
                        # Circle/Oval (check first to avoid conflict with hexagon)
                        if auto_type in [MSO_SHAPE.OVAL, 9, 10]:
                            return "circle"
                        
                        # Also check numeric values as fallback for new shapes
                        if auto_type in [7, 8]:  # Triangle types (ISOSCELES_TRIANGLE, RIGHT_TRIANGLE)
                            return "triangle"
                        elif auto_type == 182:  # Star 5-point
                            return "star"
                        elif auto_type == 56:  # Pentagon
                            return "pentagon"
                        elif auto_type == MSO_SHAPE.HEXAGON:  # Hexagon (use constant, not numeric 9)
                            return "hexagon"
                        
                        # Rectangle
                        elif auto_type in [MSO_SHAPE.RECTANGLE, 1, 2]:
                            return "square" if is_square else "rectangle"
                        
                        # Rounded rectangle
                        elif auto_type in [MSO_SHAPE.ROUNDED_RECTANGLE, 5, 6]:
                            return "roundedRectangle"
                        
                        # Line
                        elif auto_type in [20, 21, 22, 23, 24, 25]:
                            return "line"
                    except:
                        # Fallback if constants not available
                        pass
            
            # Freeform or other shapes - detect polygons by point count
            elif shape_type in [MSO_SHAPE_TYPE.FREEFORM, MSO_SHAPE_TYPE.PLACEHOLDER]:
                # IMPROVEMENT 2: Count geometry points for FREEFORM polygons (Canva compatibility)
                try:
                    point_count = None
                    
                    # Try python-pptx geometry API first
                    if hasattr(shape, 'geometry'):
                        geometry = shape.geometry
                        if hasattr(geometry, 'paths') and geometry.paths:
                            # Count points in paths
                            total_points = 0
                            for path in geometry.paths:
                                if hasattr(path, 'points'):
                                    total_points += len(path.points)
                            if total_points > 0:
                                point_count = total_points
                    
                    # Fallback: Try XML element parsing
                    if point_count is None and hasattr(shape, 'element'):
                        import xml.etree.ElementTree as ET
                        # Count pathLst/path elements and their segments
                        for path_elem in shape.element.iter():
                            tag_name = path_elem.tag.split('}')[-1] if '}' in path_elem.tag else path_elem.tag
                            if tag_name == 'pathLst' or tag_name == 'path':
                                # Count children as approximate point count
                                children = list(path_elem)
                                if len(children) > 0:
                                    point_count = len(children)
                                    break
                    
                    # Map point count to shape type
                    if point_count == 3:
                        return "triangle"
                    elif point_count == 5:
                        return "pentagon"
                    elif point_count == 6:
                        return "hexagon"
                    elif point_count == 10:
                        return "star"  # 5 outer + 5 inner points
                except Exception:
                    pass
                
                # Fallback: determine by aspect ratio
                if is_square:
                    return "square"
                else:
                    return "rectangle"
            
            # Group shapes - skip
            elif shape_type == MSO_SHAPE_TYPE.GROUP:
                return None
        
        # Fallback: determine by dimensions and properties
        if width > 0 and height > 0:
            # Very thin shape is likely a line
            if width < 1000 or height < 1000:
                return "line"
            # Square dimensions
            elif is_square:
                return "circle"  # Default to circle for square shapes
            # Rectangle
            else:
                return "rectangle"
    
    except Exception as e:
        # If detection fails, try dimension-based fallback
        try:
            if hasattr(shape, 'width') and hasattr(shape, 'height'):
                width = shape.width
                height = shape.height
                if abs(width - height) < 1000:
                    return "circle"
                else:
                    return "rectangle"
        except:
            pass
    
    return None


def normalize_hex(hex_val: str) -> Optional[str]:
    """
    Normalize hex color value to #rrggbb format (exact match from Colab code).
    """
    if not hex_val:
        return None
    
    import re
    hex_val = hex_val.strip().lower()
    
    # Match 6 hex digits
    if re.fullmatch(r'[0-9a-f]{6}', hex_val):
        return '#' + hex_val
    # Match 8 hex digits (ARGB) - use last 6
    if re.fullmatch(r'[0-9a-f]{8}', hex_val):
        return '#' + hex_val[-6:]
    
    return None


def apply_tint(hex_color: Optional[str], tint_val: int) -> Optional[str]:
    """
    Apply tint modifier (exact match from Colab code).
    OOXML tint val is in 1/100000 units (e.g., 40000 = 40%).
    Formula: new = orig + (255 - orig) * t
    """
    if hex_color is None or tint_val is None:
        return hex_color
    
    try:
        t = tint_val / 100000.0
        r = int(hex_color[1:3], 16)
        g = int(hex_color[3:5], 16)
        b = int(hex_color[5:7], 16)
        nr = int(round(r + (255 - r) * t))
        ng = int(round(g + (255 - g) * t))
        nb = int(round(b + (255 - b) * t))
        return f'#{nr:02x}{ng:02x}{nb:02x}'
    except (ValueError, TypeError):
        return hex_color


def apply_shade(hex_color: Optional[str], shade_val: int) -> Optional[str]:
    """
    Apply shade modifier (exact match from Colab code).
    Formula: new = orig * (1 - s), where s = shade_val/100000
    """
    if hex_color is None or shade_val is None:
        return hex_color
    
    try:
        s = shade_val / 100000.0
        r = int(hex_color[1:3], 16)
        g = int(hex_color[3:5], 16)
        b = int(hex_color[5:7], 16)
        nr = int(round(r * (1 - s)))
        ng = int(round(g * (1 - s)))
        nb = int(round(b * (1 - s)))
        return f'#{nr:02x}{ng:02x}{nb:02x}'
    except (ValueError, TypeError):
        return hex_color


def find_srgb_in_solidfill(elem) -> Optional[str]:
    """
    Find srgbClr in solidFill element (exact match from Colab code).
    """
    if elem is None:
        return None
    
    try:
        # Look for solidFill element (namespace-agnostic like Colab)
        for child in elem.iter():
            tag_name = child.tag.split('}')[-1] if '}' in child.tag else child.tag
            if tag_name == 'solidFill':
                # Look for srgbClr inside solidFill
                for color_elem in child.iter():
                    color_tag = color_elem.tag.split('}')[-1] if '}' in color_elem.tag else color_elem.tag
                    if color_tag == 'srgbClr':
                        val = color_elem.get('val')
                        if val:
                            return normalize_hex(val)
    except (AttributeError, TypeError):
        pass
    
    return None


def find_scheme_in_solidfill(elem) -> Optional[Dict[str, Any]]:
    """
    Find schemeClr in solidFill and extract name, tint, shade, lumMod (exact match from Colab code).
    """
    if elem is None:
        return None
    
    try:
        # Look for schemeClr (namespace-agnostic like Colab)
        sc = None
        for child in elem.iter():
            tag_name = child.tag.split('}')[-1] if '}' in child.tag else child.tag
            if tag_name == 'schemeClr':
                sc = child
                break
        
        if sc is None:
            return None
        
        name = sc.get('val')
        if not name:
            return None
        
        result = {'name': name, 'tint': None, 'shade': None, 'lumMod': None}
        
        # Extract modifiers (exact Colab approach)
        for child in list(sc):
            tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
            val_attr = child.get('val')
            
            if tag == 'tint' and val_attr:
                try:
                    result['tint'] = int(val_attr)
                except (ValueError, TypeError):
                    pass
            elif tag == 'shade' and val_attr:
                try:
                    result['shade'] = int(val_attr)
                except (ValueError, TypeError):
                    pass
            elif tag == 'lumMod' and val_attr:
                try:
                    result['lumMod'] = int(val_attr)
                except (ValueError, TypeError):
                    pass
        
        return result
    except (AttributeError, TypeError):
        return None


def extract_shape_fill_color(shape, pptx_path: Optional[str] = None) -> Optional[str]:
    """
    Extract shape fill color using exact Colab code approach.
    Works directly with shape's spPr XML element (exact match from Colab).
    """
    try:
        # Get shape's XML element (like Colab code: sp = shape element)
        if not hasattr(shape, '_element'):
            return None

        sp_elem = shape._element
        spPr = None

        # Find spPr element
        for child in sp_elem:
            tag_name = child.tag.split('}')[-1] if '}' in child.tag else child.tag
            if tag_name == 'spPr':
                spPr = child
                break

        if spPr is None:
            return None

        # Try direct srgb in solidFill
        fill_srgb = find_srgb_in_solidfill(spPr)
        if fill_srgb:
            return fill_srgb

        # If theme scheme present, extract scheme info
        scheme_info = None
        if pptx_path:
            scheme_info = find_scheme_in_solidfill(spPr)

        if scheme_info and scheme_info.get('name'):
            from converter.utils.background_extractor import get_theme_scheme_mapping
            theme_mapping = get_theme_scheme_mapping(pptx_path)
            base = theme_mapping.get(scheme_info['name']) if theme_mapping else None

            # If mapping found, apply tint/shade/lumMod if present
            resolved = base
            if not resolved:
                # fallback for common tokens
                name_low = scheme_info['name'].lower()
                if name_low in ("lt1", "bg1", "lt2", "bg2"):
                    resolved = "#ffffff"
                elif name_low in ("dk1", "tx1", "dk2", "tx2"):
                    resolved = "#000000"

            if resolved:
                # Apply tint/shade if provided
                if scheme_info.get('tint'):
                    resolved = apply_tint(resolved, scheme_info['tint'])
                if scheme_info.get('shade'):
                    resolved = apply_shade(resolved, scheme_info['shade'])
                # lumMod is more complex — optionally implement if needed
                return resolved

        return None
    except (AttributeError, TypeError, Exception):
        return None


def extract_shape_border_color(shape) -> Optional[str]:
    """
    Extract ONLY solid RGB border/stroke color from shape.
    Ignore theme colors.
    """
    try:
        if hasattr(shape, "line") and shape.line.color:
            fc = shape.line.color
            
            # MUST be RGB type
            if getattr(fc, "type", None) == 1:  # MSO_COLOR_TYPE.RGB
                if hasattr(fc, "rgb") and fc.rgb is not None:
                    return rgb_to_hex(fc.rgb)
    
    except Exception:
        pass
    
    return None


def extract_shape_border_width(shape) -> Optional[float]:
    """Extract border/stroke width from shape in points"""
    try:
        if hasattr(shape, 'line'):
            line = shape.line
            
            if hasattr(line, 'width'):
                width_emu = line.width
                if width_emu:
                    return emu_to_points(width_emu)
    
    except Exception:
        pass
    
    return 0


def treat_freeform_as_rectangle(shape, pptx_path: Optional[str] = None) -> Optional[ShapeElement]:
    """
    Treat FREEFORM shape as rectangle (for Canva compatibility).
    
    Args:
        shape: PowerPoint shape object (FREEFORM type)
        pptx_path: Path to PPTX file (for theme color resolution)
    
    Returns:
        ShapeElement as rectangle or None
    """
    try:
        # Get position and size
        left_emu = shape.left
        top_emu = shape.top
        width_emu = shape.width
        height_emu = shape.height
        
        # Convert to points and round to nearest integer
        x = round(emu_to_points(left_emu))
        y = round(emu_to_points(top_emu))
        width = round(emu_to_points(width_emu))
        height = round(emu_to_points(height_emu))
        
        # Get rotation
        rotation = 0
        if hasattr(shape, 'rotation') and shape.rotation is not None:
            rotation = int(shape.rotation / 60000)
        
        # Extract colors (with theme support)
        fill_color = extract_shape_fill_color(shape, pptx_path)
        border_color = extract_shape_border_color(shape)
        border_width = extract_shape_border_width(shape)
        
        # Format border color: use "transparent" if None (matching sample JSON format)
        if border_color is None:
            border_color = "transparent"
        
        # Create rectangle element
        element = ShapeElement(
            id=str(uuid.uuid4()),
            shapeType="rectangle",
            x=x,
            y=y,
            width=width,
            height=height,
            fillColor=fill_color,
            borderColor=border_color,
            borderWidth=round(border_width) if border_width else 0,
            rotation=rotation
        )
        
        return element
    except Exception:
        return None


def extract_shape_from_shape(shape, skip_if_has_text: bool = True, pptx_path: Optional[str] = None) -> Optional[ShapeElement]:
    """
    Extract shape element from a single PowerPoint shape.
    Returns ShapeElement if it's a valid shape, None otherwise.
    Now supports FREEFORM shapes (Canva compatibility).
    
    Args:
        shape: PowerPoint shape object
        skip_if_has_text: If True, skip shapes that have text content (default: True)
        pptx_path: Path to PPTX file (for theme color resolution)
    """
    # Skip if shape has text (those are handled by text extractor)
    if skip_if_has_text:
        if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
            # Only skip if it actually has text content
            if shape.has_text_frame and shape.text_frame.text.strip():
                return None
    
    # Get shape type
    shape_type = get_shape_type(shape)
    
    # FIX 2: Canva uses FREEFORM shapes - treat as rectangle
    if shape_type is None and hasattr(shape, 'shape_type'):
        if shape.shape_type == MSO_SHAPE_TYPE.FREEFORM:
            return treat_freeform_as_rectangle(shape, pptx_path)
    
    if shape_type is None:
        return None
    
    # Get position and size
    left_emu = shape.left
    top_emu = shape.top
    width_emu = shape.width
    height_emu = shape.height
    
    # Convert to points and round to nearest integer
    x = round(emu_to_points(left_emu))
    y = round(emu_to_points(top_emu))
    width = round(emu_to_points(width_emu))
    height = round(emu_to_points(height_emu))
    
    # Get rotation (in degrees, PowerPoint uses 60000ths of a degree)
    rotation = 0
    if hasattr(shape, 'rotation') and shape.rotation is not None:
        rotation = int(shape.rotation / 60000)
    
    # Extract colors (with theme support)
    fill_color = extract_shape_fill_color(shape, pptx_path)
    border_color = extract_shape_border_color(shape)
    border_width = extract_shape_border_width(shape)
    
    # Format border color: use "transparent" if None (matching sample JSON format)
    if border_color is None:
        border_color = "transparent"
    
    # Create shape element
    element = ShapeElement(
        id=str(uuid.uuid4()),
        shapeType=shape_type,
        x=x,
        y=y,
        width=width,
        height=height,
        fillColor=fill_color,
        borderColor=border_color,
        borderWidth=round(border_width) if border_width else 0,
        rotation=rotation
    )
    
    return element


def extract_shapes_from_pptx(pptx_path: str) -> List[List[Dict[str, Any]]]:
    """
    Extract all shape elements from a PPTX file.
    Returns a list of lists - one list per slide containing shape dictionaries.
    """
    prs = Presentation(pptx_path)
    all_slides_shapes = []
    
    for slide_idx, slide in enumerate(prs.slides):
        slide_shapes = []
        
        # Iterate through all shapes on the slide
        for shape in slide.shapes:
            # Extract shape from this shape
            shape_element = extract_shape_from_shape(shape)
            if shape_element:
                slide_shapes.append(shape_element.model_dump())
        
        all_slides_shapes.append(slide_shapes)
    
    return all_slides_shapes

