"""
Phase 3: Image Extraction from PPTX
Extracts images from slides and converts them to base64 encoded strings.
"""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import uuid
import base64
from typing import List, Dict, Any, Optional
from io import BytesIO
from converter.schemas.slide_schema import ImageElement


def emu_to_points(emu: int) -> float:
    """
    Convert EMU (English Metric Units) to points.
    1 inch = 914400 EMU = 72 points
    """
    return (emu / 914400) * 72


def image_to_base64(image_bytes: bytes, image_format: str = "png") -> str:
    """
    Convert image bytes to base64 data URL string.
    
    Args:
        image_bytes: Raw image bytes
        image_format: Image format (png, jpeg, jpg)
    
    Returns:
        Base64 data URL string (e.g., "data:image/png;base64,...")
    """
    # Normalize format
    if image_format.lower() in ['jpg', 'jpeg']:
        mime_type = "image/jpeg"
    elif image_format.lower() == 'png':
        mime_type = "image/png"
    elif image_format.lower() == 'gif':
        mime_type = "image/gif"
    else:
        mime_type = "image/png"  # Default
    
    # Encode to base64
    base64_str = base64.b64encode(image_bytes).decode('utf-8')
    
    # Return as data URL
    return f"data:{mime_type};base64,{base64_str}"


def extract_blip_image_from_graphic_frame(shape) -> Optional[ImageElement]:
    """
    Extract image from GRAPHIC_FRAME shape (used by Canva).
    Looks for embedded blip images in graphic frames.
    
    Args:
        shape: PowerPoint shape object (GRAPHIC_FRAME type)
    
    Returns:
        ImageElement or None
    """
    try:
        # Check if shape has element (XML access)
        if not hasattr(shape, 'element'):
            return None
        
        # Try to find blip element using XPath or direct XML parsing
        # Namespace: http://schemas.openxmlformats.org/drawingml/2006/main
        r_id = None
        
        try:
            # Method 1: Try lxml XPath (if available)
            try:
                from lxml import etree
                # Look for a:blip element
                blip_elements = shape.element.xpath('.//a:blip', namespaces={
                    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
                })
                
                if blip_elements:
                    blip = blip_elements[0]
                    r_id = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
            except ImportError:
                # Method 2: Fallback to direct XML element search
                import xml.etree.ElementTree as ET
                # Search for blip elements in the XML tree
                for elem in shape.element.iter():
                    # Use namespace-safe parsing
                    tag_name = elem.tag.split('}')[-1] if '}' in elem.tag else elem.tag
                    if tag_name == 'blip' or 'blip' in tag_name.lower():
                        r_id = elem.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                        if r_id:
                            break
            except Exception:
                pass
            
            if not r_id:
                return None
            
            # Get the related part (image)
            if not hasattr(shape, 'part') or not hasattr(shape.part, 'related_parts'):
                return None
            
            if r_id not in shape.part.related_parts:
                return None
            
            image_part = shape.part.related_parts[r_id]
            image_bytes = image_part.blob
            
            # Determine image format
            image_format = "png"  # Default
            if hasattr(image_part, 'content_type'):
                content_type = image_part.content_type
                if 'jpeg' in content_type or 'jpg' in content_type:
                    image_format = "jpeg"
                elif 'png' in content_type:
                    image_format = "png"
                elif 'gif' in content_type:
                    image_format = "gif"
            
            # Convert to base64
            base64_image = image_to_base64(image_bytes, image_format)
            
            # Get position and size
            left_emu = shape.left
            top_emu = shape.top
            width_emu = shape.width
            height_emu = shape.height
            
            # Convert to points and round to nearest integer
            x = round(emu_to_points(left_emu))
            y = round(emu_to_points(top_emu))
            width = round(emu_to_points(width_emu))
            height = round(emu_to_points(height_emu))
            
            # Get rotation
            rotation = 0
            if hasattr(shape, 'rotation') and shape.rotation is not None:
                rotation = int(shape.rotation / 60000)
            
            # Create image element
            element = ImageElement(
                id=str(uuid.uuid4()),
                x=x,
                y=y,
                width=width,
                height=height,
                src=base64_image,
                rotation=rotation
            )
            
            return element
        
        except ImportError:
            # lxml not available, try alternative method
            pass
        except Exception:
            pass
    
    except Exception:
        pass
    
    return None


def extract_image_from_shape(shape) -> Optional[ImageElement]:
    """
    Extract image element from a single PowerPoint shape.
    Returns ImageElement if it's a valid image, None otherwise.
    Supports: PICTURE, GRAPHIC_FRAME (Canva), and picture fills.
    """
    try:
        # Check if shape is an image
        if hasattr(shape, 'shape_type'):
            shape_type = shape.shape_type
            
            # Check if it's a picture/image
            if shape_type == MSO_SHAPE_TYPE.PICTURE:
                # Get image data
                image = shape.image
                
                # Get image bytes
                image_bytes = image.blob
                
                # Determine image format from content type or extension
                image_format = "png"  # Default
                if hasattr(image, 'ext'):
                    ext = image.ext.lower()
                    if ext in ['jpg', 'jpeg']:
                        image_format = "jpeg"
                    elif ext == 'png':
                        image_format = "png"
                    elif ext == 'gif':
                        image_format = "gif"
                
                # Convert to base64
                base64_image = image_to_base64(image_bytes, image_format)
                
                # Get position and size
                left_emu = shape.left
                top_emu = shape.top
                width_emu = shape.width
                height_emu = shape.height
                
                # Convert to points and round to nearest integer
                x = round(emu_to_points(left_emu))
                y = round(emu_to_points(top_emu))
                width = round(emu_to_points(width_emu))
                height = round(emu_to_points(height_emu))
                
                # Get rotation (in degrees, PowerPoint uses 60000ths of a degree)
                rotation = 0
                if hasattr(shape, 'rotation') and shape.rotation is not None:
                    rotation = int(shape.rotation / 60000)
                
                # Create image element
                element = ImageElement(
                    id=str(uuid.uuid4()),
                    x=x,
                    y=y,
                    width=width,
                    height=height,
                    src=base64_image,
                    rotation=rotation
                )
                
                return element
            
            # FIX 1: Canva uses GRAPHIC_FRAME for images
            # Check if it's a graphic frame (Canva images)
            elif shape_type == MSO_SHAPE_TYPE.GRAPHIC_FRAME:
                blip_image = extract_blip_image_from_graphic_frame(shape)
                if blip_image:
                    return blip_image
            
            # FIX 3: Check if shape has an image fill (for all shape types, not just AUTO_SHAPE)
            # Canva uses picture fills in freeform shapes, groups, etc.
            if hasattr(shape, 'fill'):
                fill = shape.fill
                if hasattr(fill, 'type'):
                    from pptx.enum.dml import MSO_FILL_TYPE
                    if fill.type == MSO_FILL_TYPE.PICTURE:
                            # This is a shape with picture fill
                            # Extract the picture
                            if hasattr(fill, 'picture'):
                                picture = fill.picture
                                if hasattr(picture, 'image'):
                                    image = picture.image
                                    image_bytes = image.blob
                                    
                                    # Get image format
                                    image_format = "png"
                                    if hasattr(image, 'ext'):
                                        ext = image.ext.lower()
                                        if ext in ['jpg', 'jpeg']:
                                            image_format = "jpeg"
                                        elif ext == 'png':
                                            image_format = "png"
                                    
                                    # Convert to base64
                                    base64_image = image_to_base64(image_bytes, image_format)
                                    
                                    # Get position and size
                                    left_emu = shape.left
                                    top_emu = shape.top
                                    width_emu = shape.width
                                    height_emu = shape.height
                                    
                                    # Convert to points
                                    x = emu_to_points(left_emu)
                                    y = emu_to_points(top_emu)
                                    width = emu_to_points(width_emu)
                                    height = emu_to_points(height_emu)
                                    
                                    # Get rotation
                                    rotation = 0
                                    if hasattr(shape, 'rotation') and shape.rotation is not None:
                                        rotation = int(shape.rotation / 60000)
                                    
                                    # Create image element
                                    element = ImageElement(
                                        id=str(uuid.uuid4()),
                                        x=x,
                                        y=y,
                                        width=width,
                                        height=height,
                                        src=base64_image,
                                        rotation=rotation
                                    )
                                    
                                    return element
        
        # FIX 3 (continued): Check picture fill for ANY shape type (outside shape_type check)
        # This handles picture fills in shapes that don't match PICTURE or GRAPHIC_FRAME types
        if hasattr(shape, 'fill'):
            fill = shape.fill
            if hasattr(fill, 'type'):
                from pptx.enum.dml import MSO_FILL_TYPE
                if fill.type == MSO_FILL_TYPE.PICTURE:
                    # This is a shape with picture fill
                    # Extract the picture
                    if hasattr(fill, 'picture'):
                        picture = fill.picture
                        if hasattr(picture, 'image'):
                            image = picture.image
                            image_bytes = image.blob
                            
                            # Get image format
                            image_format = "png"
                            if hasattr(image, 'ext'):
                                ext = image.ext.lower()
                                if ext in ['jpg', 'jpeg']:
                                    image_format = "jpeg"
                                elif ext == 'png':
                                    image_format = "png"
                            
                            # Convert to base64
                            base64_image = image_to_base64(image_bytes, image_format)
                            
                            # Get position and size
                            left_emu = shape.left
                            top_emu = shape.top
                            width_emu = shape.width
                            height_emu = shape.height
                            
                            # Convert to points
                            x = emu_to_points(left_emu)
                            y = emu_to_points(top_emu)
                            width = emu_to_points(width_emu)
                            height = emu_to_points(height_emu)
                            
                            # Get rotation
                            rotation = 0
                            if hasattr(shape, 'rotation') and shape.rotation is not None:
                                rotation = int(shape.rotation / 60000)
                            
                            # Create image element
                            element = ImageElement(
                                id=str(uuid.uuid4()),
                                x=x,
                                y=y,
                                width=width,
                                height=height,
                                src=base64_image,
                                rotation=rotation
                            )
                            
                            return element
    
    except Exception as e:
        # If extraction fails, return None
        pass
    
    return None


def extract_images_from_pptx(pptx_path: str) -> List[List[Dict[str, Any]]]:
    """
    Extract all image elements from a PPTX file.
    Returns a list of lists - one list per slide containing image dictionaries.
    """
    prs = Presentation(pptx_path)
    all_slides_images = []
    
    for slide_idx, slide in enumerate(prs.slides):
        slide_images = []
        
        # Iterate through all shapes on the slide
        for shape in slide.shapes:
            # Extract image from this shape
            image_element = extract_image_from_shape(shape)
            if image_element:
                slide_images.append(image_element.model_dump())
        
        all_slides_images.append(slide_images)
    
    return all_slides_images

