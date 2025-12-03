"""
Phase 1: Text Extraction from PPTX
Extracts text shapes with coordinates, font properties, alignment, rotation, and colors.
"""
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import uuid
from typing import List, Dict, Any, Optional
from converter.schemas.slide_schema import TextElement
from converter.utils.shape_extractor import extract_shape_from_shape
from converter.utils.image_extractor import extract_image_from_shape
from converter.utils.background_extractor import extract_slide_background
from converter.utils.table_extractor import extract_table_from_shape
from pptx.enum.shapes import MSO_SHAPE_TYPE


def emu_to_points(emu: int) -> float:
    """
    Convert EMU (English Metric Units) to points.
    1 inch = 914400 EMU = 72 points
    """
    return (emu / 914400) * 72


def get_alignment(pp_alignment) -> str:
    """Convert PowerPoint alignment to string"""
    if pp_alignment == PP_ALIGN.LEFT:
        return "left"
    elif pp_alignment == PP_ALIGN.CENTER:
        return "center"
    elif pp_alignment == PP_ALIGN.RIGHT:
        return "right"
    elif pp_alignment == PP_ALIGN.JUSTIFY:
        return "justify"
    else:
        return "left"  # Default


def rgb_to_hex(rgb) -> str:
    """Convert RGB color to hex string"""
    if rgb is None:
        return "#000000"  # Default black
    
    # Handle different color formats
    if hasattr(rgb, 'rgb'):
        rgb = rgb.rgb
    
    if rgb is None:
        return "#000000"
    
    # Extract RGB values
    if isinstance(rgb, int):
        # Assume it's a 32-bit integer: 0x00RRGGBB
        r = (rgb >> 16) & 0xFF
        g = (rgb >> 8) & 0xFF
        b = rgb & 0xFF
    elif hasattr(rgb, '__iter__') and len(rgb) >= 3:
        r, g, b = rgb[0], rgb[1], rgb[2]
    else:
        return "#000000"
    
    return f"#{r:02x}{g:02x}{b:02x}"


def extract_text_from_shape(shape, pptx_path: Optional[str] = None) -> List[TextElement]:
    """
    Extract text elements from a single shape.
    Returns a list of TextElement objects (one per paragraph or text frame).
    """
    elements = []
    
    if not shape.has_text_frame:
        return elements
    
    text_frame = shape.text_frame
    
    # Get shape position and size
    left_emu = shape.left
    top_emu = shape.top
    width_emu = shape.width
    height_emu = shape.height
    
    # Convert to points and round to nearest integer
    x = round(emu_to_points(left_emu))
    y = round(emu_to_points(top_emu))
    width = round(emu_to_points(width_emu))
    height = round(emu_to_points(height_emu))
    
    # Get rotation (in degrees, PowerPoint uses 60000ths of a degree)
    rotation = 0
    if hasattr(shape, 'rotation') and shape.rotation is not None:
        rotation = int(shape.rotation / 60000)
    
    # Extract text from all paragraphs
    full_text = ""
    font_properties = {
        'fontSize': 12.0,  # Default
        'fontFamily': 'Arial',  # Default
        'fontWeight': 'normal',
        'fontStyle': 'normal',
        'textDecoration': 'none',
        'color': None,
        'textAlign': 'left'
    }
    
    # Process paragraphs
    for paragraph in text_frame.paragraphs:
        para_text = paragraph.text.strip()
        if not para_text:
            continue
        
        # Get paragraph-level properties
        if paragraph.runs:
            first_run = paragraph.runs[0]
            font = first_run.font
            
            # Font size (round to nearest integer)
            if font.size:
                font_properties['fontSize'] = round(font.size.pt)
            else:
                font_properties['fontSize'] = 12
            
            # Font family
            if font.name:
                font_properties['fontFamily'] = font.name
            else:
                font_properties['fontFamily'] = 'Arial'
            
            # Font weight (bold)
            font_properties['fontWeight'] = 'bold' if (font.bold is not None and font.bold) else 'normal'
            
            # Font style (italic)
            font_properties['fontStyle'] = 'italic' if (font.italic is not None and font.italic) else 'normal'
            
            # Text decoration (underline, strikethrough)
            if hasattr(font, 'underline') and font.underline:
                font_properties['textDecoration'] = 'underline'
            elif hasattr(font, 'strike') and font.strike:
                font_properties['textDecoration'] = 'line-through'
            else:
                font_properties['textDecoration'] = 'none'
            
            # Color - accept solid RGB, but also resolve theme white/black fallback
            try:
                # Prefer direct RGB if available
                text_color = None
                if font.color and getattr(font.color, "type", None) == 1:  # MSO_COLOR_TYPE.RGB
                    if hasattr(font.color, 'rgb') and font.color.rgb is not None:
                        text_color = rgb_to_hex(font.color.rgb)

                # If no direct RGB, attempt XML theme resolution (but only fallback common tokens)
                if not text_color and hasattr(font, '_element') and pptx_path:
                    # Search for srgbClr or schemeClr in run/font XML (like table extractor does)
                    for color_elem in font._element.iter():
                        color_tag = color_elem.tag.split('}')[-1] if '}' in color_elem.tag else color_elem.tag
                        if color_tag == 'srgbClr':
                            val = color_elem.get('val')
                            if val:
                                hex_val = val.lower()
                                if not hex_val.startswith('#'):
                                    hex_val = '#' + hex_val
                                text_color = hex_val
                                break
                        elif color_tag == 'schemeClr':
                            scheme = color_elem.get('val')
                            if scheme:
                                from converter.utils.background_extractor import get_theme_scheme_mapping
                                theme_mapping = get_theme_scheme_mapping(pptx_path)
                                resolved = theme_mapping.get(scheme) if theme_mapping else None
                                if not resolved:
                                    low = scheme.lower()
                                    if low in ("lt1", "bg1", "lt2", "bg2"):
                                        resolved = "#ffffff"
                                    elif low in ("dk1", "tx1", "dk2", "tx2"):
                                        resolved = "#000000"
                                if resolved:
                                    text_color = resolved
                                    break

                font_properties['color'] = text_color if text_color else None
            except (AttributeError, TypeError):
                font_properties['color'] = None
        
        # Text alignment
        if paragraph.alignment:
            font_properties['textAlign'] = get_alignment(paragraph.alignment)
        else:
            font_properties['textAlign'] = 'left'
        
        # Combine text from all runs in paragraph
        para_full_text = ""
        for run in paragraph.runs:
            para_full_text += run.text
        
        if para_full_text.strip():
            # Create element for this paragraph
            element = TextElement(
                id=str(uuid.uuid4()),
                content=para_full_text,
                x=x,
                y=y,
                width=width,
                height=height,
                fontSize=font_properties['fontSize'],
                fontFamily=font_properties['fontFamily'],
                fontWeight=font_properties['fontWeight'],
                fontStyle=font_properties['fontStyle'],
                textDecoration=font_properties['textDecoration'],
                textAlign=font_properties['textAlign'],
                color=font_properties['color'],
                rotation=rotation
            )
            elements.append(element)
    
    # If no paragraphs but shape has text, extract as single element
    if not elements and text_frame.text.strip():
        # Try to get default font properties from shape
        if text_frame.paragraphs:
            para = text_frame.paragraphs[0]
            if para.runs:
                run = para.runs[0]
                font = run.font
                
                font_properties['fontSize'] = round(font.size.pt) if font.size else 12
                font_properties['fontFamily'] = font.name if font.name else 'Arial'
                font_properties['fontWeight'] = 'bold' if (font.bold is not None and font.bold) else 'normal'
                font_properties['fontStyle'] = 'italic' if (font.italic is not None and font.italic) else 'normal'
                
                # Text decoration
                if hasattr(font, 'underline') and font.underline:
                    font_properties['textDecoration'] = 'underline'
                elif hasattr(font, 'strike') and font.strike:
                    font_properties['textDecoration'] = 'line-through'
                else:
                    font_properties['textDecoration'] = 'none'
                
                # Color - accept solid RGB, but also resolve theme white/black fallback
                try:
                    # Prefer direct RGB if available
                    text_color = None
                    if font.color and getattr(font.color, "type", None) == 1:  # MSO_COLOR_TYPE.RGB
                        if hasattr(font.color, 'rgb') and font.color.rgb is not None:
                            text_color = rgb_to_hex(font.color.rgb)

                    # If no direct RGB, attempt XML theme resolution (but only fallback common tokens)
                    if not text_color and hasattr(font, '_element') and pptx_path:
                        # Search for srgbClr or schemeClr in run/font XML (like table extractor does)
                        for color_elem in font._element.iter():
                            color_tag = color_elem.tag.split('}')[-1] if '}' in color_elem.tag else color_elem.tag
                            if color_tag == 'srgbClr':
                                val = color_elem.get('val')
                                if val:
                                    hex_val = val.lower()
                                    if not hex_val.startswith('#'):
                                        hex_val = '#' + hex_val
                                    text_color = hex_val
                                    break
                            elif color_tag == 'schemeClr':
                                scheme = color_elem.get('val')
                                if scheme:
                                    from converter.utils.background_extractor import get_theme_scheme_mapping
                                    theme_mapping = get_theme_scheme_mapping(pptx_path)
                                    resolved = theme_mapping.get(scheme) if theme_mapping else None
                                    if not resolved:
                                        low = scheme.lower()
                                        if low in ("lt1", "bg1", "lt2", "bg2"):
                                            resolved = "#ffffff"
                                        elif low in ("dk1", "tx1", "dk2", "tx2"):
                                            resolved = "#000000"
                                    if resolved:
                                        text_color = resolved
                                        break

                    font_properties['color'] = text_color if text_color else None
                except (AttributeError, TypeError):
                    font_properties['color'] = None
            
            # Alignment
            if para.alignment:
                font_properties['textAlign'] = get_alignment(para.alignment)
            else:
                font_properties['textAlign'] = 'left'
        
        element = TextElement(
            id=str(uuid.uuid4()),
            content=text_frame.text,
            x=x,
            y=y,
            width=width,
            height=height,
            fontSize=font_properties['fontSize'],
            fontFamily=font_properties['fontFamily'],
            fontWeight=font_properties['fontWeight'],
            fontStyle=font_properties['fontStyle'],
            textDecoration=font_properties['textDecoration'],
            textAlign=font_properties['textAlign'],
            color=font_properties['color'],
            rotation=rotation
        )
        elements.append(element)
    
    return elements


def extract_text_from_pptx(pptx_path: str) -> List[Dict[str, Any]]:
    """
    Extract all text elements from a PPTX file.
    Returns a list of slides, each containing text elements.
    All coordinates are scaled from PowerPoint points to Presentera canvas (1024×576).
    """
    from converter.utils.scaling import get_slide_dimensions, calculate_scale_factor, scale_element_coordinates
    
    prs = Presentation(pptx_path)
    
    # Get slide dimensions and calculate scale factors
    ppt_width, ppt_height = get_slide_dimensions(prs)
    scale_x, scale_y = calculate_scale_factor(ppt_width, ppt_height)
    
    slides_data = []
    
    for slide_idx, slide in enumerate(prs.slides):
        slide_elements = []
        
        # Iterate through all shapes on the slide
        for shape_idx, shape in enumerate(slide.shapes):
            # Check if shape is a table first (tables should be extracted separately)
            if hasattr(shape, 'has_table') and shape.has_table:
                table_element = extract_table_from_shape(shape, pptx_path=pptx_path)
                if table_element:
                    table_dict = table_element.model_dump()
                    # Scale table element coordinates
                    table_dict = scale_element_coordinates(table_dict, scale_x, scale_y)
                    # Also scale cell dimensions and round to nearest integer
                    if 'cellWidth' in table_dict:
                        table_dict['cellWidth'] = round(table_dict['cellWidth'] * scale_x)
                    if 'cellHeight' in table_dict:
                        table_dict['cellHeight'] = round(table_dict['cellHeight'] * scale_y)
                    # Scale cell font sizes and border widths and round to nearest integer
                    if 'data' in table_dict:
                        for row in table_dict['data']:
                            for cell in row:
                                if 'fontSize' in cell:
                                    cell['fontSize'] = round(cell['fontSize'] * scale_x)
                                if 'borderWidth' in cell:
                                    cell['borderWidth'] = round(cell['borderWidth'] * scale_x)
                    slide_elements.append(table_dict)
                    continue  # Skip other extraction for table shapes
            
            # Extract image (images are separate from text/shapes)
            image_element = extract_image_from_shape(shape)
            if image_element:
                image_dict = image_element.model_dump()
                # Scale image element coordinates
                image_dict = scale_element_coordinates(image_dict, scale_x, scale_y)
                slide_elements.append(image_dict)
                continue  # Skip text/shape extraction for image shapes
            
            # Extract text from this shape
            text_elements = extract_text_from_shape(shape, pptx_path=pptx_path)
            # Convert TextElement objects to dictionaries and scale
            for text_elem in text_elements:
                text_dict = text_elem.model_dump()
                # Scale text element coordinates
                text_dict = scale_element_coordinates(text_dict, scale_x, scale_y)
                slide_elements.append(text_dict)
            
            # Extract shape (only if it doesn't have text content and is not a chart)
            # This ensures we don't duplicate text boxes as shapes and don't extract chart shapes as rectangles
            has_text = hasattr(shape, 'has_text_frame') and shape.has_text_frame and shape.text_frame.text.strip()
            # Safely check if shape is a chart using helper function
            from converter.utils.chart_extractor import is_chart_shape
            is_chart = is_chart_shape(shape)
            
            if not has_text and not is_chart:
                # FIX 5: Handle GROUP shapes (Canva) - process children
                if hasattr(shape, 'shape_type') and shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    if hasattr(shape, 'shapes'):
                        for child_shape in shape.shapes:
                            # Skip chart children - they'll be extracted separately
                            child_is_chart = is_chart_shape(child_shape)
                            if child_is_chart:
                                continue
                            child_shape_elem = extract_shape_from_shape(child_shape, skip_if_has_text=True, pptx_path=pptx_path)
                            if child_shape_elem:
                                child_dict = child_shape_elem.model_dump()
                                # Scale shape element coordinates
                                child_dict = scale_element_coordinates(child_dict, scale_x, scale_y)
                                slide_elements.append(child_dict)
                else:
                    shape_element = extract_shape_from_shape(shape, skip_if_has_text=True, pptx_path=pptx_path)
                    if shape_element:
                        shape_dict = shape_element.model_dump()
                        # Scale shape element coordinates
                        shape_dict = scale_element_coordinates(shape_dict, scale_x, scale_y)
                        slide_elements.append(shape_dict)
        
        # Extract charts from slide (separate extractor)
        from converter.utils.chart_extractor import extract_charts_from_slide
        chart_elements = extract_charts_from_slide(slide, pptx_path, scale_x, scale_y)
        slide_elements.extend(chart_elements)
        
        # Extract slide background properties
        background_info = extract_slide_background(slide, pptx_path=pptx_path, slide_index=slide_idx)
        
        # --- BACKGROUND HANDLING ---
        # NO background image elements are created for any background type
        # Only backgroundColor and backgroundImage properties are set in slide_data
        # No extra images are added to elements array
        
        # Create slide data with all required fields
        slide_data = {
            "id": str(uuid.uuid4()),
            "elements": slide_elements,
            "thumbnail": None,
            # Legacy fields for backward compatibility
            "backgroundColor": background_info.get("backgroundColor"),
            "backgroundImage": background_info.get("backgroundImage"),
            "backgroundSize": background_info.get("backgroundSize"),
            "backgroundPosition": background_info.get("backgroundPosition"),
            "backgroundRepeat": background_info.get("backgroundRepeat")
        }
        slides_data.append(slide_data)
    
    return slides_data

