# converter/utils/chart_extractor.py
import zipfile
import io
import uuid
from lxml import etree
from openpyxl import load_workbook
from openpyxl.utils import range_boundaries
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Use your project's scaling helper
from converter.utils.scaling import emu_to_points

NS = {
    "c": "http://schemas.openxmlformats.org/drawingml/2006/chart",
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
}

DEFAULT_PALETTE = [
    "#10b981",  # teal/green
    "#0ea5e9",  # blue
    "#8b5cf6",  # purple
    "#f59e0b",  # amber
    "#f43f5e",  # red/pink
    "#06b6d4",  # cyan
    "#84cc16",  # lime
]


def detect_chart_type(chart_root):
    """Return 'bar'|'line'|'pie' or 'unknown'"""
    if chart_root.find(".//c:barChart", NS) is not None:
        return "bar"
    if chart_root.find(".//c:lineChart", NS) is not None:
        return "line"
    if chart_root.find(".//c:pieChart", NS) is not None or chart_root.find(".//c:doughnutChart", NS) is not None:
        return "pie"
    return "unknown"


def get_chart_title(chart_root):
    """Extract chart title text (if any)"""
    t = chart_root.find(".//c:title//a:t", NS)
    return t.text if t is not None else "chart name"


def _read_cached_categories(chart_root):
    cats = []
    for pt in chart_root.findall(".//c:cat//c:strCache//c:pt", NS):
        v = pt.find("c:v", NS)
        if v is not None and v.text is not None:
            cats.append(v.text)
    return cats


def _read_cached_series(chart_root):
    series = {}
    for ser in chart_root.findall(".//c:ser", NS):
        # series name
        name_el = ser.find(".//c:tx//c:v", NS)
        name = name_el.text if name_el is not None else "Series"
        vals = []
        for pt in ser.findall(".//c:val//c:numCache//c:pt", NS):
            v = pt.find("c:v", NS)
            if v is not None and v.text is not None:
                try:
                    vals.append(float(v.text))
                except Exception:
                    # leave as-is if not numeric
                    vals.append(v.text)
        series[name] = vals
    return series


def _extract_series_colors(chart_root):
    """
    Try to extract series color from chart XML.
    Returns list of color strings (hex without '#') per series (may be shorter than series count).
    """
    colors = []
    # Look for series spPr / solidFill / srgbClr val
    for ser in chart_root.findall(".//c:ser", NS):
        clr = None
        sppr = ser.find(".//c:spPr", NS)
        if sppr is not None:
            # try a:solidFill/a:srgbClr/@val
            srgb = sppr.find(".//a:solidFill//a:srgbClr", NS)
            if srgb is not None and 'val' in srgb.attrib:
                clr = srgb.attrib['val']
            else:
                # try schemeClr or prstClr -> fallback
                scheme = sppr.find(".//a:solidFill//a:schemeClr", NS)
                if scheme is not None and 'val' in scheme.attrib:
                    # scheme val e.g. 'accent1' - we can't resolve theme here easily, so skip
                    clr = None
        if clr:
            # convert RRGGBB to #rrggbb
            colors.append("#" + clr.lower())
        else:
            colors.append(None)
    return colors


def _read_embedded_excel_from_rels(chart_path, zipf):
    """
    Check chart rels for an embedded workbook and return (sheet_obj or None)
    Note: returns workbook object (openpyxl) or None.
    """
    # chart_path example: 'ppt/charts/chart1.xml'
    rels_path = chart_path.replace("charts/chart", "charts/_rels/chart") + ".rels"
    try:
        rels_xml = zipf.read(rels_path)
    except KeyError:
        return None

    rels_root = etree.fromstring(rels_xml)
    for rel in rels_root.findall("r:Relationship", NS):
        target = rel.get("Target")
        if not target:
            continue
        if "embeddings" in target:
            # normalize path
            emb_path = "ppt/" + target.replace("../", "")
            try:
                wb_bytes = zipf.read(emb_path)
                wb = load_workbook(io.BytesIO(wb_bytes), data_only=True)
                return wb
            except Exception:
                return None
    return None


def _load_series_from_excel(chart_root, wb):
    """
    Given chart_root xml and an openpyxl workbook, read ranges referenced by each <c:ser>.
    Returns categories list and series dict.
    """
    sheet = wb.active
    cats = []
    series = {}

    for ser in chart_root.findall(".//c:ser", NS):
        # name
        name_el = ser.find(".//c:tx//c:v", NS)
        series_name = name_el.text if name_el is not None else "Series"

        # category reference (c:cat/c:f)
        cat_ref = ser.find(".//c:cat//c:f", NS)
        if cat_ref is not None and not cats:
            try:
                min_col, min_row, max_col, max_row = range_boundaries(cat_ref.text)
                for row in sheet.iter_rows(min_row=min_row, max_row=max_row, min_col=min_col, max_col=max_col):
                    cats.append(row[0].value)
            except Exception:
                pass

        # value reference (c:val/c:f)
        vals = []
        val_ref = ser.find(".//c:val//c:f", NS)
        if val_ref is not None:
            try:
                min_col, min_row, max_col, max_row = range_boundaries(val_ref.text)
                for row in sheet.iter_rows(min_row=min_row, max_row=max_row, min_col=min_col, max_col=max_col):
                    vals.append(row[0].value)
            except Exception:
                pass

        series[series_name] = vals

    return cats, series


def create_chart_element(shape, chart_root, categories, series_data, series_colors, scale_x=1, scale_y=1):
    """
    Build the element dict that matches sample_8.json structure.
    """
    chart_type = detect_chart_type(chart_root)
    chart_title = get_chart_title(chart_root)

    # Geometry (scale from EMU -> points -> canvas)
    x = round(emu_to_points(shape.left) * scale_x)
    y = round(emu_to_points(shape.top) * scale_y)
    width = round(emu_to_points(shape.width) * scale_x)
    height = round(emu_to_points(shape.height) * scale_y)
    rotation = int(shape.rotation / 60000) if hasattr(shape, 'rotation') and shape.rotation else 0

    element = {
        "id": str(uuid.uuid4()),
        "type": "chart",
        "chartType": chart_type,
        "x": x,
        "y": y,
        "width": width,
        "height": height,
        "chartName": chart_title,
        "backgroundColor": "#ffffff",
        "rotation": rotation
    }

    # Build content based on type
    if chart_type in ("bar", "line", "unknown"):
        # Use categories and series[] structure (sample_8.json)
        element["labels"] = categories or []
        element["showXAxis"] = True
        # showYAxis guess: line charts usually have y axis, bar may or may not
        element["showYAxis"] = True if chart_type == "line" else False

        formatted_series = []
        palette = DEFAULT_PALETTE

        # If series_colors provided as list
        for i, (sname, svals) in enumerate(series_data.items()):
            # compute barColors: try to use extracted color else palette
            color = None
            if i < len(series_colors) and series_colors[i]:
                color = series_colors[i]
            else:
                color = palette[i % len(palette)]

            bar_colors = [color] * max(1, len(svals))
            formatted_series.append({
                "name": sname,
                "values": [v if v is not None else 0 for v in svals],
                "barColors": bar_colors
            })
        element["series"] = formatted_series

    elif chart_type == "pie":
        # pie: sample_8.json stores labels + values + barColors + color (primary)
        # In chart XML, pie often has single series; take the first series values
        element["labels"] = categories or []
        # flatten first series values
        first_series_vals = []
        if series_data:
            # take first series
            first_series_vals = list(next(iter(series_data.values())))
        element["values"] = [v if v is not None else 0 for v in first_series_vals]
        # colors: map series_colors or palette
        colors = []
        for i in range(len(element["values"])):
            if i < len(series_colors) and series_colors[i]:
                colors.append(series_colors[i])
            else:
                colors.append(DEFAULT_PALETTE[i % len(DEFAULT_PALETTE)])
        element["barColors"] = colors
        # primary color (match sample)
        element["color"] = colors[0] if colors else DEFAULT_PALETTE[0]

    return element


def is_chart_shape(shape):
    """
    Safely check if a shape is a chart shape.
    Returns True if shape is a chart, False otherwise.
    """
    try:
        # First check shape_type (safe, doesn't raise exceptions)
        if hasattr(shape, 'shape_type') and shape.shape_type == MSO_SHAPE_TYPE.CHART:
            return True
        # Try to access chart property (may raise ValueError if not a chart)
        try:
            _ = shape.chart
            return True
        except (ValueError, AttributeError):
            return False
    except Exception:
        return False


def extract_charts_from_slide(slide, pptx_path, scale_x=1, scale_y=1):
    """
    Extract chart elements from a slide object.
    Returns list of element dicts (matching sample_8.json).
    Args:
        slide: python-pptx slide object
        pptx_path: path to the PPTX file (we need zip access to read embedded workbooks and chart xml)
        scale_x, scale_y: canvas scaling factors (use your scaling.calculate_scale_factor results)
    """
    chart_elements = []

    # open zip once
    try:
        z = zipfile.ZipFile(pptx_path, "r")
    except Exception:
        z = None

    for shape in slide.shapes:
        # skip non-chart shapes
        if not is_chart_shape(shape):
            # If group, still need to check children
            if hasattr(shape, 'shape_type') and shape.shape_type == MSO_SHAPE_TYPE.GROUP and hasattr(shape, 'shapes'):
                for child in shape.shapes:
                    if is_chart_shape(child):
                        # treat child as chart
                        try:
                            chart_part_name = child.chart.part.partname.lstrip("/")
                        except Exception:
                            # fallback to continue (can't map)
                            continue

                        try:
                            chart_root = etree.fromstring(z.read(chart_part_name))
                        except Exception:
                            continue

                        # Try embedded excel
                        wb = _read_embedded_excel_from_rels(chart_part_name, z)
                        if wb:
                            cats, series = _load_series_from_excel(chart_root, wb)
                        else:
                            cats = _read_cached_categories(chart_root)
                            series = _read_cached_series(chart_root)

                        series_colors = _extract_series_colors(chart_root)
                        elem = create_chart_element(child, chart_root, cats, series, series_colors, scale_x, scale_y)
                        chart_elements.append(elem)
                continue
            else:
                continue

        # for top-level chart shapes
        try:
            chart_part_name = shape.chart.part.partname.lstrip("/")
        except Exception:
            # sometimes shape may not expose chart attribute
            continue

        try:
            chart_root = etree.fromstring(z.read(chart_part_name))
        except Exception:
            continue

        # load embedded excel if present
        wb = _read_embedded_excel_from_rels(chart_part_name, z)
        if wb:
            cats, series = _load_series_from_excel(chart_root, wb)
        else:
            cats = _read_cached_categories(chart_root)
            series = _read_cached_series(chart_root)

        series_colors = _extract_series_colors(chart_root)
        elem = create_chart_element(shape, chart_root, cats, series, series_colors, scale_x, scale_y)
        chart_elements.append(elem)

    # close zip if opened here
    if z:
        z.close()

    return chart_elements
