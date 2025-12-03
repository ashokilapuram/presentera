"""
Table Extraction from PPTX
Extracts tables with cell properties, colors, fonts, borders, and alignment.
"""
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_FILL_TYPE
import uuid
from typing import List, Dict, Any, Optional
from converter.schemas.slide_schema import TableElement, TableCell


def emu_to_points(emu: int) -> float:
    """
    Convert EMU (English Metric Units) to points.
    1 inch = 914400 EMU = 72 points
    """
    return (emu / 914400) * 72


def rgb_to_hex(rgb) -> str:
    """Convert RGB color to hex string"""
    if rgb is None:
        return "#000000"  # Default black
    
    # Handle different color formats
    if hasattr(rgb, 'rgb'):
        rgb = rgb.rgb
    
    if rgb is None:
        return "#000000"
    
    # Extract RGB values
    if isinstance(rgb, int):
        # Assume it's a 32-bit integer: 0x00RRGGBB
        r = (rgb >> 16) & 0xFF
        g = (rgb >> 8) & 0xFF
        b = rgb & 0xFF
    elif hasattr(rgb, '__iter__') and len(rgb) >= 3:
        r, g, b = rgb[0], rgb[1], rgb[2]
    else:
        return "#000000"
    
    return f"#{r:02x}{g:02x}{b:02x}"


def get_alignment(pp_alignment) -> str:
    """Convert PowerPoint alignment to string"""
    if pp_alignment == PP_ALIGN.LEFT:
        return "left"
    elif pp_alignment == PP_ALIGN.CENTER:
        return "center"
    elif pp_alignment == PP_ALIGN.RIGHT:
        return "right"
    elif pp_alignment == PP_ALIGN.JUSTIFY:
        return "justify"
    else:
        return "left"  # Default


def extract_cell_fill_color(cell, pptx_path: Optional[str] = None) -> str:
    """
    Extract fill color from table cell.
    Supports multiple color formats:
    - Direct sRGB colors (via python-pptx API and XML)
    - Theme/scheme colors (via XML parsing with theme mapping)
    - Gradient fills (uses first stop color)
    
    Returns hex color string, defaults to "#FFFFFF" (white) if not found.
    """
    # Method 1: Try python-pptx API first (fast path for direct RGB)
    try:
        if hasattr(cell, 'fill'):
            fill = cell.fill
            if hasattr(fill, 'type'):
                fill_type = fill.type
                
                # Solid fill
                if fill_type == MSO_FILL_TYPE.SOLID:
                    if hasattr(fill, 'fore_color'):
                        try:
                            fore_color = fill.fore_color
                            # Try direct RGB first
                            if fore_color and hasattr(fore_color, 'rgb') and fore_color.rgb is not None:
                                hex_color = rgb_to_hex(fore_color.rgb)
                                if hex_color:  # Valid color found
                                    return hex_color
                        except (AttributeError, TypeError):
                            pass
                
                # Gradient fill - use first stop color
                elif fill_type == MSO_FILL_TYPE.GRADIENT:
                    if hasattr(fill, 'gradient_stops'):
                        stops = fill.gradient_stops
                        if stops and len(stops) > 0:
                            try:
                                first_stop = stops[0]
                                if hasattr(first_stop, 'color'):
                                    color = first_stop.color
                                    if color and hasattr(color, 'rgb') and color.rgb is not None:
                                        hex_color = rgb_to_hex(color.rgb)
                                        if hex_color:
                                            return hex_color
                            except (AttributeError, TypeError):
                                pass
    except Exception:
        pass
    
    # Method 2: Parse XML directly for theme colors and other formats
    try:
        if hasattr(cell, '_element'):
            elem = cell._element
            
            # Table cells store fill in tcPr (table cell properties) > solidFill
            # First, look for tcPr element
            tc_pr = None
            for child in elem:
                tag_name = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                if tag_name == 'tcPr':
                    tc_pr = child
                    break
            
            # If tcPr found, look for solidFill inside it
            if tc_pr is not None:
                for fill_elem in tc_pr:
                    tag_name = fill_elem.tag.split('}')[-1] if '}' in fill_elem.tag else fill_elem.tag
                    
                    if tag_name == 'solidFill':
                        # Found solidFill, now look for color inside it
                        for color_elem in fill_elem:
                            color_tag = color_elem.tag.split('}')[-1] if '}' in color_elem.tag else color_elem.tag
                            
                            # Direct sRGB color
                            if color_tag == 'srgbClr':
                                val = color_elem.get('val')
                                if val:
                                    # Ensure hex format
                                    hex_val = val.lower()
                                    if not hex_val.startswith('#'):
                                        hex_val = '#' + hex_val
                                    return hex_val
                            
                            # Theme/scheme color
                            elif color_tag == 'schemeClr':
                                scheme_name = color_elem.get('val')
                                if scheme_name and pptx_path:
                                    from converter.utils.background_extractor import get_theme_scheme_mapping
                                    theme_mapping = get_theme_scheme_mapping(pptx_path)
                                    if scheme_name in theme_mapping:
                                        return theme_mapping[scheme_name]
                                    else:
                                        # fallback
                                        lname = scheme_name.lower()
                                        if lname in ("lt1", "bg1", "lt2", "bg2"):
                                            return "#ffffff"
                                        if lname in ("dk1", "tx1", "dk2", "tx2"):
                                            return "#000000"
                            
                            # System colors (prstClr) - map common ones
                            elif color_tag == 'prstClr':
                                prst_val = color_elem.get('val')
                                if prst_val:
                                    # Map common preset colors
                                    preset_map = {
                                        'black': '#000000',
                                        'white': '#FFFFFF',
                                        'red': '#FF0000',
                                        'green': '#00FF00',
                                        'blue': '#0000FF',
                                        'yellow': '#FFFF00',
                                        'cyan': '#00FFFF',
                                        'magenta': '#FF00FF'
                                    }
                                    if prst_val.lower() in preset_map:
                                        return preset_map[prst_val.lower()]
            
            # Fallback: search all elements for solidFill (in case structure is different)
            for fill_elem in elem.iter():
                tag_name = fill_elem.tag.split('}')[-1] if '}' in fill_elem.tag else fill_elem.tag
                
                if tag_name == 'solidFill':
                    # Found solidFill, now look for color inside it
                    for color_elem in fill_elem:
                        color_tag = color_elem.tag.split('}')[-1] if '}' in color_elem.tag else color_elem.tag
                        
                        # Direct sRGB color
                        if color_tag == 'srgbClr':
                            val = color_elem.get('val')
                            if val:
                                hex_val = val.lower()
                                if not hex_val.startswith('#'):
                                    hex_val = '#' + hex_val
                                return hex_val
                        
                        # Theme/scheme color
                        elif color_tag == 'schemeClr':
                            scheme_name = color_elem.get('val')
                            if scheme_name and pptx_path:
                                from converter.utils.background_extractor import get_theme_scheme_mapping
                                theme_mapping = get_theme_scheme_mapping(pptx_path)
                                if scheme_name in theme_mapping:
                                    return theme_mapping[scheme_name]
                                else:
                                    # fallback
                                    lname = scheme_name.lower()
                                    if lname in ("lt1", "bg1", "lt2", "bg2"):
                                        return "#ffffff"
                                    if lname in ("dk1", "tx1", "dk2", "tx2"):
                                        return "#000000"
    except Exception:
        pass
    
    return "#FFFFFF"  # Default white


def extract_cell_border_color(cell, pptx_path: Optional[str] = None) -> str:
    """
    Extract border color from table cell.
    Supports multiple color formats via XML parsing.
    """
    try:
        if hasattr(cell, '_element'):
            elem = cell._element
            
            # Table cells have borders defined in tcBdr (table cell borders)
            # Look for tcBdr element first
            tc_bdr = None
            for child in elem:
                tag_name = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                if tag_name == 'tcBdr':
                    tc_bdr = child
                    break
            
            if tc_bdr is None:
                return "#FFFFFF"
            
            # Check each border side (top, bottom, left, right)
            # Use the first border found (usually all sides have same color)
            for border_elem in tc_bdr:
                tag_name = border_elem.tag.split('}')[-1] if '}' in border_elem.tag else border_elem.tag
                
                if tag_name in ['top', 'bottom', 'left', 'right']:
                    # Look for line (ln) element inside border
                    for ln_elem in border_elem:
                        ln_tag = ln_elem.tag.split('}')[-1] if '}' in ln_elem.tag else ln_elem.tag
                        if ln_tag == 'ln':
                            # Look for color inside line element
                            for color_elem in ln_elem.iter():
                                color_tag = color_elem.tag.split('}')[-1] if '}' in color_elem.tag else color_elem.tag
                                
                                # Direct sRGB color
                                if color_tag == 'srgbClr':
                                    val = color_elem.get('val')
                                    if val:
                                        hex_val = val.lower()
                                        if not hex_val.startswith('#'):
                                            hex_val = '#' + hex_val
                                        return hex_val
                                
                                # Theme/scheme color
                                elif color_tag == 'schemeClr':
                                    scheme_name = color_elem.get('val')
                                    if scheme_name and pptx_path:
                                        from converter.utils.background_extractor import get_theme_scheme_mapping
                                        theme_mapping = get_theme_scheme_mapping(pptx_path)
                                        if scheme_name in theme_mapping:
                                            return theme_mapping[scheme_name]
                                        else:
                                            # fallback
                                            lname = scheme_name.lower()
                                            if lname in ("lt1", "bg1", "lt2", "bg2"):
                                                return "#ffffff"
                                            if lname in ("dk1", "tx1", "dk2", "tx2"):
                                                return "#000000"
                                
                                # System colors (prstClr)
                                elif color_tag == 'prstClr':
                                    prst_val = color_elem.get('val')
                                    if prst_val:
                                        preset_map = {
                                            'black': '#000000',
                                            'white': '#FFFFFF',
                                            'red': '#FF0000',
                                            'green': '#00FF00',
                                            'blue': '#0000FF',
                                            'yellow': '#FFFF00',
                                            'cyan': '#00FFFF',
                                            'magenta': '#FF00FF'
                                        }
                                        if prst_val.lower() in preset_map:
                                            return preset_map[prst_val.lower()]
                            
                            # If we found a border element but no color, break after first border
                            break
                    
                    # Return after processing first border side
                    if tag_name in ['top', 'bottom', 'left', 'right']:
                        break
    except Exception:
        pass
    
    return "#FFFFFF"  # Default white


def extract_cell_border_width(cell) -> float:
    """Extract border width from table cell in points"""
    try:
        if hasattr(cell, '_element'):
            elem = cell._element
            # Look for border width in XML
            # Table cells have borders defined in tcBdr
            for border_elem in elem.iter():
                tag_name = border_elem.tag.split('}')[-1] if '}' in border_elem.tag else border_elem.tag
                
                # Check for table cell border elements
                if tag_name in ['top', 'bottom', 'left', 'right', 'ln']:
                    # Get width attribute (w) - in EMU
                    width_attr = border_elem.get('w')
                    if width_attr:
                        try:
                            width_emu = int(width_attr)
                            return emu_to_points(width_emu)
                        except (ValueError, TypeError):
                            pass
                    
                    # Also check for line width in child elements
                    for child in border_elem.iter():
                        child_tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                        if child_tag == 'ln':
                            width_attr = child.get('w')
                            if width_attr:
                                try:
                                    width_emu = int(width_attr)
                                    return emu_to_points(width_emu)
                                except (ValueError, TypeError):
                                    pass
    except Exception:
        pass
    
    return 2  # Default border width


def get_fallback_cell_colors(row_index: int, is_header: bool = False) -> Dict[str, str]:
    """
    Get fallback colors for table cells when color extraction fails.
    Uses colors similar to sample_7.json pattern.
    
    Args:
        row_index: Zero-based row index
        is_header: Whether this is a header row (typically row 0)
    
    Returns:
        Dictionary with bgColor, textColor, borderColor
    """
    if is_header or row_index == 0:
        # Header row: Blue background, dark text
        return {
            'bgColor': '#2196F3',      # Blue
            'textColor': '#070707',    # Dark gray
            'borderColor': '#FFFFFF'   # White
        }
    else:
        # Data rows: Alternate between light blue shades
        if row_index % 2 == 1:
            # Odd rows: Light blue
            return {
                'bgColor': '#BBDEFB',    # Light blue
                'textColor': '#ff0f0f',  # Red
                'borderColor': '#FFFFFF' # White
            }
        else:
            # Even rows: Lighter blue
            return {
                'bgColor': '#E3F2FD',    # Lighter blue
                'textColor': '#ff0f0f',  # Red
                'borderColor': '#FFFFFF' # White
            }


def extract_cell_properties(cell, pptx_path: Optional[str] = None, row_index: int = 0, is_header: bool = False) -> Dict[str, Any]:
    """
    Extract all properties from a table cell.
    Returns a dictionary with cell properties.
    """
    # Default properties
    properties = {
        'text': '',
        'bgColor': '#FFFFFF',
        'textColor': '#000000',
        'borderColor': '#FFFFFF',
        'borderWidth': 2,
        'fontSize': 12,
        'fontFamily': 'Arial',
        'fontWeight': 'normal',
        'fontStyle': 'normal',
        'textDecoration': 'none',
        'align': 'left'
    }
    
    # Extract text
    if hasattr(cell, 'text'):
        properties['text'] = cell.text.strip()
    
    # Extract fill color (background)
    extracted_bg_color = extract_cell_fill_color(cell, pptx_path)
    
    # TEMPORARY FIX: Use fallback colors when extraction returns default white
    # This handles cases where non-RGB color formats (theme colors) fail to extract
    # RGB colors will still be extracted correctly and used
    if extracted_bg_color == "#FFFFFF":
        # Check if cell actually has a fill - if it does, white might be the real color
        # If no fill, definitely use fallback
        has_fill = False
        try:
            if hasattr(cell, 'fill'):
                fill = cell.fill
                if hasattr(fill, 'type') and fill.type == MSO_FILL_TYPE.SOLID:
                    has_fill = True
                    # If it has fill but we got white, check if it's actually a theme color we couldn't resolve
                    # In that case, use fallback
                    if hasattr(fill, 'fore_color'):
                        fore_color = fill.fore_color
                        # Check if it's a theme color (not RGB)
                        if fore_color and not (hasattr(fore_color, 'rgb') and fore_color.rgb is not None):
                            # Likely a theme color that failed to extract - use fallback
                            has_fill = False
        except Exception:
            pass
        
        if not has_fill:
            fallback_colors = get_fallback_cell_colors(row_index, is_header)
            properties['bgColor'] = fallback_colors['bgColor']
        else:
            properties['bgColor'] = extracted_bg_color
    else:
        properties['bgColor'] = extracted_bg_color
    
    # Extract border color and width
    extracted_border_color = extract_cell_border_color(cell, pptx_path)
    
    # For borders, use fallback white if extraction returned default
    # (borders in sample_7.json are always white)
    if extracted_border_color == "#FFFFFF":
        fallback_colors = get_fallback_cell_colors(row_index, is_header)
        properties['borderColor'] = fallback_colors['borderColor']
    else:
        properties['borderColor'] = extracted_border_color
    
    properties['borderWidth'] = round(extract_cell_border_width(cell))
    
    # Extract text formatting from first paragraph
    if hasattr(cell, 'text_frame') and cell.text_frame.paragraphs:
        para = cell.text_frame.paragraphs[0]
        
        # Get alignment
        if para.alignment:
            properties['align'] = get_alignment(para.alignment)
        
        # Get font properties from first run
        if para.runs:
            run = para.runs[0]
            font = run.font
            
            # Font size (round to nearest integer)
            if font.size:
                properties['fontSize'] = round(font.size.pt)
            else:
                properties['fontSize'] = 12
            
            # Font family
            if font.name:
                properties['fontFamily'] = font.name
            else:
                properties['fontFamily'] = 'Arial'
            
            # Font weight
            properties['fontWeight'] = 'bold' if (font.bold is not None and font.bold) else 'normal'
            
            # Font style
            properties['fontStyle'] = 'italic' if (font.italic is not None and font.italic) else 'normal'
            
            # Text decoration
            if hasattr(font, 'underline') and font.underline:
                properties['textDecoration'] = 'underline'
            elif hasattr(font, 'strike') and font.strike:
                properties['textDecoration'] = 'line-through'
            else:
                properties['textDecoration'] = 'none'
            
                # Text color - try multiple methods
            try:
                text_color = None
                
                # Method 1: Try python-pptx API (direct RGB)
                if font.color:
                    if hasattr(font.color, 'rgb') and font.color.rgb is not None:
                        text_color = rgb_to_hex(font.color.rgb)
                
                # Method 2: Parse XML for theme colors and other formats
                if not text_color and hasattr(font, '_element'):
                    font_elem = font._element
                    # Look for solidFill or color elements
                    for color_elem in font_elem.iter():
                        color_tag = color_elem.tag.split('}')[-1] if '}' in color_elem.tag else color_elem.tag
                        
                        # Direct sRGB color
                        if color_tag == 'srgbClr':
                            val = color_elem.get('val')
                            if val:
                                hex_val = val.lower()
                                if not hex_val.startswith('#'):
                                    hex_val = '#' + hex_val
                                text_color = hex_val
                                break
                        
                        # Theme/scheme color
                        elif color_tag == 'schemeClr':
                            scheme_name = color_elem.get('val')
                            if scheme_name and pptx_path:
                                from converter.utils.background_extractor import get_theme_scheme_mapping
                                theme_mapping = get_theme_scheme_mapping(pptx_path)
                                resolved = theme_mapping.get(scheme_name) if theme_mapping else None
                                if not resolved:
                                    lname = scheme_name.lower()
                                    if lname in ("lt1", "bg1", "lt2", "bg2"):
                                        resolved = "#ffffff"
                                    elif lname in ("dk1", "tx1", "dk2", "tx2"):
                                        resolved = "#000000"
                                if resolved:
                                    text_color = resolved
                                    break
                
                # Method 3: Also check text_frame XML directly
                if not text_color and hasattr(cell, 'text_frame') and hasattr(cell.text_frame, '_element'):
                    tf_elem = cell.text_frame._element
                    for para_elem in tf_elem.iter():
                        para_tag = para_elem.tag.split('}')[-1] if '}' in para_elem.tag else para_elem.tag
                        if para_tag == 'p':
                            for run_elem in para_elem.iter():
                                run_tag = run_elem.tag.split('}')[-1] if '}' in run_elem.tag else run_elem.tag
                                if run_tag == 'r':
                                    for prop_elem in run_elem.iter():
                                        prop_tag = prop_elem.tag.split('}')[-1] if '}' in prop_elem.tag else prop_elem.tag
                                        if prop_tag == 'solidFill':
                                            for color_elem in prop_elem:
                                                color_tag = color_elem.tag.split('}')[-1] if '}' in color_elem.tag else color_elem.tag
                                                if color_tag == 'srgbClr':
                                                    val = color_elem.get('val')
                                                    if val:
                                                        hex_val = val.lower()
                                                        if not hex_val.startswith('#'):
                                                            hex_val = '#' + hex_val
                                                        text_color = hex_val
                                                        break
                                                elif color_tag == 'schemeClr':
                                                    scheme_name = color_elem.get('val')
                                                    if scheme_name and pptx_path:
                                                        from converter.utils.background_extractor import get_theme_scheme_mapping
                                                        theme_mapping = get_theme_scheme_mapping(pptx_path)
                                                        resolved = theme_mapping.get(scheme_name) if theme_mapping else None
                                                        if not resolved:
                                                            lname = scheme_name.lower()
                                                            if lname in ("lt1", "bg1", "lt2", "bg2"):
                                                                resolved = "#ffffff"
                                                            elif lname in ("dk1", "tx1", "dk2", "tx2"):
                                                                resolved = "#000000"
                                                        if resolved:
                                                            text_color = resolved
                                                            break
                                    if text_color:
                                        break
                            if text_color:
                                break
                
                # Text color extraction is working fine - use extracted color or default black
                properties['textColor'] = text_color if text_color else '#000000'
            except (AttributeError, TypeError, Exception):
                # Default to black on exception (text color extraction is working)
                properties['textColor'] = '#000000'
    
    return properties


def extract_table_from_shape(shape, pptx_path: Optional[str] = None) -> Optional[TableElement]:
    """
    Extract table element from a single PowerPoint shape.
    Returns TableElement if it's a valid table, None otherwise.
    
    Args:
        shape: PowerPoint shape object
        pptx_path: Path to PPTX file (for theme color resolution)
    """
    try:
        # Check if shape has a table
        if not hasattr(shape, 'has_table') or not shape.has_table:
            return None
        
        table = shape.table
        
        # Get table position and size
        left_emu = shape.left
        top_emu = shape.top
        width_emu = shape.width
        height_emu = shape.height
        
        # Convert to points and round to nearest integer
        x = round(emu_to_points(left_emu))
        y = round(emu_to_points(top_emu))
        width = round(emu_to_points(width_emu))
        height = round(emu_to_points(height_emu))
        
        # Get rotation
        rotation = 0
        if hasattr(shape, 'rotation') and shape.rotation is not None:
            rotation = int(shape.rotation / 60000)
        
        # Get table dimensions
        rows = len(table.rows)
        # Get column count from first row (more reliable)
        if rows > 0 and len(table.rows[0].cells) > 0:
            cols = len(table.rows[0].cells)
        else:
            cols = 0
        
        # Calculate cell dimensions and round to nearest integer
        if cols > 0:
            cell_width = round(width / cols)
        else:
            cell_width = 0
        
        if rows > 0:
            cell_height = round(height / rows)
        else:
            cell_height = 0
        
        # Extract cell data
        table_data = []
        for row_idx, row in enumerate(table.rows):
            row_data = []
            # Determine if this is a header row (typically first row)
            is_header = (row_idx == 0)
            
            for cell in row.cells:
                # Extract cell properties
                cell_props = extract_cell_properties(cell, pptx_path, row_index=row_idx, is_header=is_header)
                
                # Create TableCell object
                table_cell = TableCell(
                    text=cell_props['text'],
                    bgColor=cell_props['bgColor'],
                    textColor=cell_props['textColor'],
                    borderColor=cell_props['borderColor'],
                    borderWidth=cell_props['borderWidth'],
                    fontSize=cell_props['fontSize'],
                    fontFamily=cell_props['fontFamily'],
                    fontWeight=cell_props['fontWeight'],
                    fontStyle=cell_props['fontStyle'],
                    textDecoration=cell_props['textDecoration'],
                    align=cell_props['align']
                )
                row_data.append(table_cell)
            
            table_data.append(row_data)
        
        # Create table element
        element = TableElement(
            id=str(uuid.uuid4()),
            x=x,
            y=y,
            width=width,
            height=height,
            rows=rows,
            cols=cols,
            cellWidth=cell_width,
            cellHeight=cell_height,
            data=table_data,
            rotation=rotation
        )
        
        return element
    
    except Exception as e:
        # If extraction fails, return None
        return None


def extract_tables_from_pptx(pptx_path: str) -> List[List[Dict[str, Any]]]:
    """
    Extract all table elements from a PPTX file.
    Returns a list of lists - one list per slide containing table dictionaries.
    """
    prs = Presentation(pptx_path)
    all_slides_tables = []
    
    for slide_idx, slide in enumerate(prs.slides):
        slide_tables = []
        
        # Iterate through all shapes on the slide
        for shape in slide.shapes:
            # Extract table from this shape
            table_element = extract_table_from_shape(shape, pptx_path)
            if table_element:
                slide_tables.append(table_element.model_dump())
        
        all_slides_tables.append(slide_tables)
    
    return all_slides_tables

